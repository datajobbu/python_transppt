'''
실행 코드: python3 ppt_translate.py [파일 경로/파일명] [파파고 API 아이디] [파파고 API 비밀번호]

#####예시#####
python3 ppt_translate.py ../python_transppt/examples/lecture_ex.pptx id pw
> translated_lecture_ex.pptx 파일 동일 위치에 생성됨
'''
import sys
from pptx import Presentation
from papago_api_pre import translate


def input_analyzer():
    try:
        full_name = sys.argv[1]
        nmt_id = sys.argv[2]
        nmt_pw = sys.argv[3]
        
        temp_list = full_name.split('/')
        file_name = temp_list.pop()
        temp_list.append("")
        path = "/".join(temp_list)

        return path, file_name, nmt_id, nmt_pw

    except:
        print("Input Error")
        sys.exit()


def ppt_translator(prs, nmt_id, nmt_pw):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                translated = translate(paragraph.text, nmt_id, nmt_pw)
                paragraph.clear()
                
                run = paragraph.add_run()
                run.text = translated


def main():
    _path, _filename, _id, _pw = input_analyzer()
    print('Ready..')
    prs = Presentation(_path + _filename)
    ppt_translator(prs, _id, _pw)
    prs.save(_path+"translated_"+_filename)
    print('Done!')
    

if __name__ == "__main__":
    main()